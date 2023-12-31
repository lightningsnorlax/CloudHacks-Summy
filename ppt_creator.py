from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests
import io
from PIL import Image
import json

import time


def add_image_slide(width, height, slide, image_path, position):
    if position == 'top_right':
        image = slide.shapes.add_picture(
            image_file=image_path,  left=width / 1.9, top=height / 10)
    elif position == 'middle_left':
        image = slide.shapes.add_picture(
            image_file=image_path,  left=width / 7, top=height / 2.9)
    elif position == 'bottom_right':
        image = slide.shapes.add_picture(
            image_file=image_path,  left=width / 1.9, top=height / 1.7)

    image.z_order = 2


def add_text_box(width, height, slide, text_title, text_body, position):
    if position == 'top_left':
        textbox = slide.shapes.add_textbox(
            left=width / 7, top=height / 10, width=Cm(8), height=Cm(1))
        text_frame = textbox.text_frame
        text_frame.text = text_title

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                font.bold = True
                font.name = "Patrick Hand"
                font.size = Pt(12.5)
                font.color.rgb = RGBColor(14, 83, 102)

        image = slide.shapes.add_picture(
            image_file="lightbulb.png",  left=width / 9, top=height / 10)
        image.z_order = 2

        textbox = slide.shapes.add_textbox(
            left=width / 7, top=height / 7, width=Cm(8), height=Cm(5))
        text_frame = textbox.text_frame
        text_frame.text = text_body

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                font.name = "Open Sans"
                font.size = Pt(11.4)
                font.color.rgb = RGBColor(14, 83, 102)

    elif position == 'middle_right':
        textbox = slide.shapes.add_textbox(
            left=width / 2, top=height / 2.9, width=Cm(8), height=Cm(1))
        text_frame = textbox.text_frame
        text_frame.text = text_title

        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.RIGHT

            for run in paragraph.runs:
                font = run.font
                font.bold = True
                font.name = "Patrick Hand"
                font.size = Pt(12.5)
                font.color.rgb = RGBColor(14, 83, 102)

        image = slide.shapes.add_picture(
            image_file="lightbulb.png",  left=width / 1.13, top=height / 2.9)
        image.z_order = 2

        textbox = slide.shapes.add_textbox(
            left=width / 2, top=height / 2.6, width=Cm(8), height=Cm(5))
        text_frame = textbox.text_frame
        text_frame.text = text_body

        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.RIGHT

            for run in paragraph.runs:
                font = run.font
                font.name = "Open Sans"
                font.size = Pt(11.4)
                font.color.rgb = RGBColor(14, 83, 102)

    elif position == 'bottom_left':
        textbox = slide.shapes.add_textbox(
            left=width / 7, top=height / 1.7, width=Cm(8), height=Cm(1))
        text_frame = textbox.text_frame
        text_frame.text = text_title

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                font.bold = True
                font.name = "Patrick Hand"
                font.size = Pt(12.5)
                font.color.rgb = RGBColor(14, 83, 102)

        image = slide.shapes.add_picture(
            image_file="lightbulb.png",  left=width / 9, top=height / 1.7)
        image.z_order = 2

        textbox = slide.shapes.add_textbox(
            left=width / 7, top=height / 1.58, width=Cm(8), height=Cm(5))
        text_frame = textbox.text_frame
        text_frame.text = text_body

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                font.name = "Open Sans"
                font.size = Pt(11.4)
                font.color.rgb = RGBColor(14, 83, 102)


def createPowerPoint(jsonObject):
    body = jsonObject["body"]
    # Create a new presentation
    prs = Presentation()

    # Set the slide width and height based on the A4 paper aspect ratio
    slide_width = Inches(210 / 25.4)  # Convert mm to inches
    slide_height = Inches(297 / 25.4)

    # Add a blank slide with the specified dimensions
    # Use layout index 5 for a blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Remove the original empty text box
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.text_frame.text == "":
            sp = shape
            slide.shapes._spTree.remove(sp._element)

    # Set the slide dimensions
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    # Add a picture to the slide as a background
    image_path = 'background.png'  # Replace with your image file path

    left = 0
    top = 0
    width = slide_width
    height = slide_height

    background_picture = slide.shapes.add_picture(
        image_path, left, top, width, height)

    # Send the background picture to the back
    background_picture.z_order = 3

    # Add a text box to the slide
    text_left = (slide_width - Cm(14.82)) / 2  # Center align horizontally
    text_top = Cm(0.43)
    text_width = Cm(14.82)
    text_height = Cm(1.47)

    textbox = slide.shapes.add_textbox(
        text_left, text_top, text_width, text_height)
    text_frame = textbox.text_frame

    # Add content to the text box
    paragraph = text_frame.add_paragraph()
    run = paragraph.add_run()
    run.text = jsonObject["title"]

    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            font = run.font
            font.bold = True
            font.name = "Patrick Hand"
            font.size = Pt(28)
            font.color.rgb = RGBColor(14, 83, 102)

    # Set font size and style
    font = run.font
    font.size = Pt(28)  # Change the font size to 28 points

    # Center align the text horizontally
    paragraph.alignment = PP_ALIGN.CENTER

    # APIS Stuff
    prompt_list = [body[0]["image"], body[1]["image"], body[2]["image"]]
    image_list = []

    API_URL = "https://api-inference.huggingface.co/models/runwayml/stable-diffusion-v1-5"

    headers = {"Authorization": "Bearer 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"}

    def query(payload):
        response = requests.post(API_URL, headers=headers, json=payload)
        return response.content

    pos = ["top_right", "middle_left", "bottom_right"]

    for i in range(0, 3):

        image_bytes = query({
            "inputs": prompt_list[i],
        })
        # time.sleep(10)
        print(image_bytes)
        image = Image.open(io.BytesIO(image_bytes))

        output_file = f"./images/output_image{i}.jpg"
        image.save(output_file)
        # time.sleep(10)
        print(pos[i])
        add_image_slide(slide, output_file, pos[i])
        # image_list.append(output_file)

        print(f"Image saved to {output_file}")

    # add_image_slide(slide, image_list[1], "middle_left")
    # add_image_slide(slide, image_list[2], "bottom_right")

    add_text_box(width, height, slide, text_title=body[0]["subheading"],
                 text_body=body[0]["text"], position='top_left')
    add_text_box(width, height, slide, text_title=body[1]["subheading"],
                 text_body=body[1]["text"], position='middle_right')
    add_text_box(width, height, slide, text_title=body[2]["subheading"],
                 text_body=body[2]["text"], position='bottom_left')

    # Save the presentation to a file
    prs.save('format.pptx')


# Opening JSON file
f = open('ideal_model_output.json')
print("hi")

hi = json.load(f)

createPowerPoint(hi)
