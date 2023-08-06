from flask import Blueprint, render_template, redirect, url_for, request, flash
from flask import json, jsonify
import random
from flask_cors import CORS, cross_origin
import base64
import requests
import numpy as np
import re
import datetime
from dotenv import load_dotenv
import os
import subprocess
import pinecone
# Langchain
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.text_splitter import CharacterTextSplitter
from langchain.vectorstores import Pinecone
from langchain.document_loaders import TextLoader
from langchain.chat_models import ChatOpenAI
from langchain.chains import RetrievalQA, LLMChain, LLMCheckerChain
from langchain.chains.question_answering import load_qa_chain
from langchain.callbacks import wandb_tracing_enabled
from langchain.prompts import (
    PromptTemplate,
    ChatPromptTemplate,
    HumanMessagePromptTemplate,
)
from langchain.prompts.few_shot import FewShotPromptTemplate

from typing import Optional
from langchain.chains import SimpleSequentialChain, SequentialChain

from langchain.chains.openai_functions import (
    create_openai_fn_chain,
    create_structured_output_chain,
)
from langchain.schema import HumanMessage, AIMessage, ChatMessage
import openai

from langchain.chat_models import ChatOpenAI
from langchain.document_loaders import WebBaseLoader
from langchain.chains.summarize import load_summarize_chain

import asyncio

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests
import io
from PIL import Image
import json

load_dotenv()

main = Blueprint('main', __name__)

OPENAI_API_KEY = os.environ.get('OPENAI_API_KEY')
PINECONE_API_KEY = os.environ.get("PINECONE_API_KEY")
INDEX_NAME = os.environ.get("PINECONE_INDEX_NAME")
PINECONE_ENVIRONMENT = os.environ.get("PINECONE_ENVIRONMENT")
HUGGING_FACE_API_KEY = os.environ.get("HUGGING_FACE_API_KEY")

openai.api_key = OPENAI_API_KEY

pinecone.init(
    api_key=PINECONE_API_KEY,
    environment=PINECONE_ENVIRONMENT,
)
embeddings = OpenAIEmbeddings(model='text-embedding-ada-002')
index = pinecone.Index(INDEX_NAME)
model_name = "gpt-3.5-turbo-0613"
temperature = 0.0
llm = ChatOpenAI(model_name=model_name, temperature=temperature)
vectorstore = Pinecone(index, embeddings.embed_query, "text")


@main.route('/')
@main.route('/index')
@main.route('/home')
def index():
    return render_template("/index.html")


@main.route('/question')
def question():
    return render_template("/question.html")


commandArr = [
    "wget -r -np -nd -A.html,.txt,.tmp -P websites",
    "python cleaner.py",
    "python chunker.py",
    "python vectorizor.py",
]


@main.route("/getSummary", methods=["GET", "POST"])
def getSummary():
    if request.method == "POST":
        if request.json['url']:
            url = request.json['url']
            print(url)
            loader = WebBaseLoader(url)
            docs = loader.load()

            chain = load_summarize_chain(llm, chain_type="refine")
            summary = chain.run(docs)
            print(summary)
            return jsonify({'message': f'OK', "summary": summary})
        return jsonify({'message': f'invalid url'})
    return render_template("index.html")


@main.route("/askQn", methods=["GET", "POST"])
def askQuestion():
    if request.method == "POST":
        if request.json['url']:
            url = request.json['url']
            print(url)
            for command in commandArr:
                if "wget" in command:
                    command = command + " " + url
                result = subprocess.run(
                    command, shell=True, capture_output=True, text=True, cwd="./Sn33k")
            print("Summary stored")
            if request.json['question']:
                qa_chain = RetrievalQA.from_chain_type(
                    llm=llm,
                    chain_type="refine",
                    retriever=vectorstore.as_retriever(),
                    # verbose=True,
                )
                query = request.json['question']
                print(query)
                chainRes = SequentialChain(chains=[qa_chain], input_variables=[
                    "query"]).run(query=query)
                return jsonify({'message': "OK", "answer": chainRes})
        return jsonify({'message': f'invalid url'})
    return render_template("index.html")


@main.route("/generatePoster", methods=["POST"])
def generatePoster():
    if request.method == "POST":
        if request.json['url']:
            url = request.json['url']
            print(url)
            loader = WebBaseLoader(url)
            docs = loader.load()
            json_schema = {
                "title": "Title of Poster",
                "description": "Poster Topic Summarised",
                "type": "object",
                "properties": {
                    "body": {
                        "type": "array",
                        "description": "List of points to be included in the summary, it is a dictionary with keys: 1. image  2. text",
                        "items": {
                            "type": "object"
                        },
                        "properties": {
                            "subheading": {
                                "type": "string",
                                "description": "Title of the text paragraph"
                            },
                            "image": {
                                "type": "string",
                                "description": "A description of the image, should be specific"
                            },
                            "text": {
                                "type": "string",
                                "description": "Informative and concisive summary on the article provided about the point"
                            }
                        },
                        "required": ["image", "text", "subheading"]
                    }
                },
                "required": ["list_of_points"]
            }

            summary_chain = load_summarize_chain(
                llm, chain_type="stuff")
            post_template = PromptTemplate(
                template="""Goal:Generate a summary poster based on custom knowledge base (Information below) and user query. Two components 1.Poster assets descriptions 2.Poster content.\n\n------------------Custom knowledge base:------------------\n{query}\n------------------End of Custom Knowledge base.------------------\nExample Format output: dict(
                "title": "Title of poster",
                "body":[
                dict(
                    "subheading": "subheading description1",
                    "image": "image description1...",
                    "text": "text description1"
                ),
                dict(
                    "subheading": "subheading description2",
                    "image": "image description2...",
                    "text": "text description2"
                ),
                dict(
                    "subheading": "subheading description3",
                    "image": "image description3...",
                    "text": "text description3"
                ),
                ]
            )\n\nUsing the above information in Custom knowledge base, generate a poster content that summarises the information""",
                input_variables=["query"],
                validate_template=False
            )
            post_chain = create_structured_output_chain(
                json_schema, llm, post_template, verbose=True)
            chainRes = SequentialChain(chains=[post_chain], input_variables=[
                                       "query"]).run(query=summary_chain.run(docs))
            return jsonify({'message': f'OK', "chain": chainRes})
    return render_template("index.html")


def add_image_slide(width, height, slide, image_path, position):
    if position == 'top_right':
        image = slide.shapes.add_picture(
            image_file=image_path,  left=width / 1.9, top=height / 10)
    elif position == 'middle_left':
        image = slide.shapes.add_picture(
            image_file=image_path,  left=width / 7, top=height / 2.9)
    elif position == 'bottom_right':
        image = slide.shapes.add_picture(
            image_file=image_path,  left=width / 1.9, top=height / 1.7)

    image.z_order = 2


def add_text_box(width, height, slide, text_title, text_body, position):
    if position == 'top_left':
        textbox = slide.shapes.add_textbox(
            left=width / 7, top=height / 10, width=Cm(8), height=Cm(1))
        text_frame = textbox.text_frame
        text_frame.text = text_title

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                font.bold = True
                font.name = "Patrick Hand"
                font.size = Pt(12.5)
                font.color.rgb = RGBColor(14, 83, 102)

        image = slide.shapes.add_picture(
            image_file="lightbulb.png",  left=width / 9, top=height / 10)
        image.z_order = 2

        textbox = slide.shapes.add_textbox(
            left=width / 7, top=height / 7, width=Cm(8), height=Cm(5))
        text_frame = textbox.text_frame
        text_frame.text = text_body

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                font.name = "Open Sans"
                font.size = Pt(11.4)
                font.color.rgb = RGBColor(14, 83, 102)

    elif position == 'middle_right':
        textbox = slide.shapes.add_textbox(
            left=width / 2, top=height / 2.9, width=Cm(8), height=Cm(1))
        text_frame = textbox.text_frame
        text_frame.text = text_title

        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.RIGHT

            for run in paragraph.runs:
                font = run.font
                font.bold = True
                font.name = "Patrick Hand"
                font.size = Pt(12.5)
                font.color.rgb = RGBColor(14, 83, 102)

        image = slide.shapes.add_picture(
            image_file="lightbulb.png",  left=width / 1.13, top=height / 2.9)
        image.z_order = 2

        textbox = slide.shapes.add_textbox(
            left=width / 2, top=height / 2.6, width=Cm(8), height=Cm(5))
        text_frame = textbox.text_frame
        text_frame.text = text_body

        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.RIGHT

            for run in paragraph.runs:
                font = run.font
                font.name = "Open Sans"
                font.size = Pt(11.4)
                font.color.rgb = RGBColor(14, 83, 102)

    elif position == 'bottom_left':
        textbox = slide.shapes.add_textbox(
            left=width / 7, top=height / 1.7, width=Cm(8), height=Cm(1))
        text_frame = textbox.text_frame
        text_frame.text = text_title

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                font.bold = True
                font.name = "Patrick Hand"
                font.size = Pt(12.5)
                font.color.rgb = RGBColor(14, 83, 102)

        image = slide.shapes.add_picture(
            image_file="lightbulb.png",  left=width / 9, top=height / 1.7)
        image.z_order = 2

        textbox = slide.shapes.add_textbox(
            left=width / 7, top=height / 1.58, width=Cm(8), height=Cm(5))
        text_frame = textbox.text_frame
        text_frame.text = text_body

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                font.name = "Open Sans"
                font.size = Pt(11.4)
                font.color.rgb = RGBColor(14, 83, 102)

